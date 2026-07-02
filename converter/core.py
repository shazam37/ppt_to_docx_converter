from __future__ import annotations

import io
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH


class ConversionError(Exception):
    pass


class PPTToDocConverter:
    def __init__(self, progress_cb=None, log_cb=None):
        """
        progress_cb(percent: int, message: str)
        log_cb(message: str)
        """
        self._progress = progress_cb or (lambda p, m: None)
        self._log = log_cb or (lambda m: None)

    def convert(self, input_path: str | Path, output_path: str | Path):
        input_path = Path(input_path)
        output_path = Path(output_path)

        if not input_path.exists():
            raise ConversionError(f"Input file not found: {input_path}")

        try:
            self._log(f"Loading presentation: {input_path.name}")
            prs = Presentation(str(input_path))
        except Exception as e:
            raise ConversionError(f"Cannot open presentation: {e}") from e

        try:
            self._log("Creating Word document")
            doc = Document()

            # Set up document margins (narrower for better space usage)
            sections = doc.sections
            for section in sections:
                section.top_margin = Inches(0.5)
                section.bottom_margin = Inches(0.5)
                section.left_margin = Inches(0.75)
                section.right_margin = Inches(0.75)

            total_slides = len(prs.slides)
            self._log(f"Processing {total_slides} slides")

            for idx, slide in enumerate(prs.slides, 1):
                self._progress(int((idx / total_slides) * 90), f"Processing slide {idx}/{total_slides}")
                self._log(f"Processing slide {idx}")

                # Add slide number heading
                heading = doc.add_heading(f"Slide {idx}", level=1)
                heading.alignment = WD_ALIGN_PARAGRAPH.LEFT

                # Process slide layout name if available
                try:
                    layout_name = slide.slide_layout.name
                    if layout_name:
                        p = doc.add_paragraph()
                        run = p.add_run(f"Layout: {layout_name}")
                        run.font.size = Pt(9)
                        run.font.italic = True
                        run.font.color.rgb = RGBColor(128, 128, 128)
                except:
                    pass

                # Process all shapes in the slide
                self._process_shapes(slide.shapes, doc)

                # Add page break after each slide except the last one
                if idx < total_slides:
                    doc.add_page_break()

            self._progress(95, "Saving document")
            self._log(f"Saving to: {output_path.name}")
            doc.save(str(output_path))

            self._progress(100, "Conversion complete")
            self._log("Conversion completed successfully")

        except Exception as e:
            raise ConversionError(f"Conversion failed: {e}") from e

    def _process_shapes(self, shapes, doc, indent_level=0):
        """Process shapes recursively (handles grouped shapes)"""
        for shape in shapes:
            try:
                # Handle grouped shapes recursively
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._log(f"  Processing group with {len(shape.shapes)} shapes")
                    self._process_shapes(shape.shapes, doc, indent_level + 1)
                    continue

                # Handle text boxes and shapes with text
                if shape.has_text_frame:
                    self._extract_text_frame(shape, doc, indent_level)

                # Handle tables
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    self._extract_table(shape, doc)

                # Handle pictures
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._extract_picture(shape, doc)

                # Handle charts (extract title if available)
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    self._extract_chart(shape, doc)

            except Exception as e:
                self._log(f"  Warning: Could not process shape: {e}")
                continue

    def _extract_text_frame(self, shape, doc, indent_level=0):
        """Extract text from a text frame"""
        text_frame = shape.text_frame

        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            # Create a paragraph in the document
            p = doc.add_paragraph()

            # Apply indentation for grouped shapes
            if indent_level > 0:
                p.paragraph_format.left_indent = Inches(0.25 * indent_level)

            # Process runs to preserve formatting
            for run in paragraph.runs:
                if not run.text.strip():
                    continue

                doc_run = p.add_run(run.text)

                # Copy formatting
                try:
                    if run.font.bold:
                        doc_run.font.bold = True
                    if run.font.italic:
                        doc_run.font.italic = True
                    if run.font.underline:
                        doc_run.font.underline = True
                    if run.font.size:
                        doc_run.font.size = Pt(min(run.font.size.pt, 14))  # Cap at 14pt
                    if run.font.color and run.font.color.type == 1:  # RGB color
                        rgb = run.font.color.rgb
                        if rgb:
                            doc_run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                except:
                    pass

            # Apply paragraph alignment
            try:
                if paragraph.alignment is not None:
                    p.alignment = paragraph.alignment
            except:
                pass

    def _extract_table(self, shape, doc):
        """Extract table data"""
        try:
            ppt_table = shape.table
            rows = len(ppt_table.rows)
            cols = len(ppt_table.columns)

            self._log(f"  Extracting table ({rows}x{cols})")

            # Create table in document
            doc_table = doc.add_table(rows=rows, cols=cols)
            doc_table.style = 'Light Grid Accent 1'

            for i, row in enumerate(ppt_table.rows):
                for j, cell in enumerate(row.cells):
                    try:
                        text = cell.text.strip()
                        doc_table.rows[i].cells[j].text = text

                        # Make first row bold (header)
                        if i == 0:
                            for paragraph in doc_table.rows[i].cells[j].paragraphs:
                                for run in paragraph.runs:
                                    run.font.bold = True
                    except:
                        continue

            # Add spacing after table
            doc.add_paragraph()

        except Exception as e:
            self._log(f"  Warning: Could not extract table: {e}")

    def _extract_picture(self, shape, doc):
        """Extract and embed picture"""
        try:
            image = shape.image
            image_bytes = image.blob

            # Get image dimensions
            width = shape.width
            height = shape.height

            # Convert EMU to inches (914400 EMU = 1 inch)
            width_inches = width / 914400
            height_inches = height / 914400

            # Cap maximum width at 6 inches (maintain aspect ratio)
            max_width = 6.0
            if width_inches > max_width:
                ratio = max_width / width_inches
                width_inches = max_width
                height_inches = height_inches * ratio

            self._log(f"  Embedding image ({width_inches:.2f}\" x {height_inches:.2f}\")")

            # Add image to document
            doc.add_picture(io.BytesIO(image_bytes), width=Inches(width_inches))
            doc.add_paragraph()  # Add spacing

        except Exception as e:
            self._log(f"  Warning: Could not extract image: {e}")
            # Add placeholder text
            p = doc.add_paragraph()
            run = p.add_run("[Image could not be extracted]")
            run.font.italic = True
            run.font.color.rgb = RGBColor(255, 0, 0)

    def _extract_chart(self, shape, doc):
        """Extract chart information"""
        try:
            if shape.has_chart:
                chart = shape.chart
                chart_title = chart.chart_title.text_frame.text if chart.has_title else "Chart"

                self._log(f"  Found chart: {chart_title}")

                p = doc.add_paragraph()
                run = p.add_run(f"[Chart: {chart_title}]")
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 0, 255)

                # Add note about chart data
                p = doc.add_paragraph()
                run = p.add_run("(Chart data not extracted - refer to original presentation)")
                run.font.size = Pt(9)
                run.font.italic = True
                run.font.color.rgb = RGBColor(128, 128, 128)

        except Exception as e:
            self._log(f"  Warning: Could not extract chart: {e}")
