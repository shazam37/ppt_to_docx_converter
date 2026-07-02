from __future__ import annotations

import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.shared import OxmlElement
from docx.oxml.ns import qn


class ConversionError(Exception):
    pass


class PPTToDocConverterClean:
    def __init__(self, progress_cb=None, log_cb=None):
        """
        Clean converter focused on visual fidelity without metadata clutter
        progress_cb(percent: int, message: str)
        log_cb(message: str)
        """
        self._progress = progress_cb or (lambda p, m: None)
        self._log = log_cb or (lambda m: None)

    def convert(self, input_path: str | Path, output_path: str | Path):
        input_path = Path(input_path)
        output_path = Path(output_path)

        if not input_path.exists():
            raise ConversionError(f"Input file not found: {input_path}")

        try:
            self._log(f"Loading presentation: {input_path.name}")
            prs = Presentation(str(input_path))
        except Exception as e:
            raise ConversionError(f"Cannot open presentation: {e}") from e

        try:
            self._log("Creating Word document")
            doc = Document()

            # Set up document to match presentation aspect ratio
            sections = doc.sections
            for section in sections:
                section.top_margin = Inches(0.5)
                section.bottom_margin = Inches(0.5)
                section.left_margin = Inches(0.75)
                section.right_margin = Inches(0.75)
                # Match presentation dimensions
                section.page_width = Inches(prs.slide_width / 914400)
                section.page_height = Inches(prs.slide_height / 914400)

            total_slides = len(prs.slides)
            self._log(f"Processing {total_slides} slides")

            for idx, slide in enumerate(prs.slides, 1):
                self._progress(int((idx / total_slides) * 90), f"Processing slide {idx}/{total_slides}")
                self._log(f"Processing slide {idx}")

                # Simple slide header
                heading = doc.add_heading(f"Slide {idx}", level=1)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

                # Add minimal separator
                p = doc.add_paragraph("─" * 80)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER

                # Sort shapes by position for natural reading order
                sorted_shapes = self._sort_shapes_by_position(slide.shapes)

                # Process shapes - content only, no metadata
                self._process_shapes(sorted_shapes, doc, slide)

                # Add page break after each slide except the last one
                if idx < total_slides:
                    doc.add_page_break()

            self._progress(95, "Saving document")
            self._log(f"Saving to: {output_path.name}")
            doc.save(str(output_path))

            self._progress(100, "Conversion complete")
            self._log("Conversion completed successfully")

        except Exception as e:
            raise ConversionError(f"Conversion failed: {e}") from e

    def _sort_shapes_by_position(self, shapes):
        """Sort shapes by vertical then horizontal position"""
        shape_list = []
        for shape in shapes:
            try:
                top = shape.top
                left = shape.left
                shape_list.append((top, left, shape))
            except:
                shape_list.append((float('inf'), float('inf'), shape))

        shape_list.sort(key=lambda x: (x[0], x[1]))
        return [s[2] for s in shape_list]

    def _process_shapes(self, shapes, doc, slide, indent_level=0):
        """Process shapes cleanly - content only"""
        for shape in shapes:
            try:
                # Skip embedded objects and connectors
                if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    continue

                # Handle grouped shapes recursively
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._log(f"  Processing group with {len(shape.shapes)} shapes")
                    sorted_group_shapes = self._sort_shapes_by_position(shape.shapes)
                    self._process_shapes(sorted_group_shapes, doc, slide, indent_level)
                    continue

                # Handle text boxes and shapes with text
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    text_content = shape.text_frame.text.strip()
                    if text_content:  # Only process if has content
                        self._extract_text_frame_clean(shape, doc, indent_level)

                # Handle tables
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    self._extract_table_clean(shape, doc)

                # Handle pictures
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._extract_picture_clean(shape, doc)

                # Handle charts - extract as image or data table
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    self._extract_chart_clean(shape, doc)

            except Exception as e:
                self._log(f"  Warning: Could not process shape: {e}")
                continue

    def _extract_text_frame_clean(self, shape, doc, indent_level=0):
        """Extract text cleanly without metadata"""
        text_frame = shape.text_frame

        for paragraph in text_frame.paragraphs:
            # Skip empty paragraphs unless they're meaningful spacing
            text_content = paragraph.text.strip()

            # Create paragraph
            p = doc.add_paragraph()

            # Apply indentation
            total_indent = indent_level
            try:
                if paragraph.level > 0:
                    total_indent += paragraph.level
            except:
                pass

            if total_indent > 0:
                p.paragraph_format.left_indent = Inches(0.25 * total_indent)

            # Process runs to preserve formatting
            has_content = False
            for run in paragraph.runs:
                if not run.text:
                    continue

                has_content = True
                doc_run = p.add_run(run.text)

                # Copy formatting
                try:
                    if run.font.bold is not None:
                        doc_run.font.bold = run.font.bold
                    if run.font.italic is not None:
                        doc_run.font.italic = run.font.italic
                    if run.font.underline is not None:
                        doc_run.font.underline = run.font.underline

                    # Preserve font size
                    if run.font.size:
                        doc_run.font.size = run.font.size

                    # Preserve font name
                    if run.font.name:
                        doc_run.font.name = run.font.name

                    # Preserve color
                    if run.font.color and run.font.color.type == 1:
                        rgb = run.font.color.rgb
                        if rgb:
                            doc_run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])

                except Exception as e:
                    pass

            # If no content, remove the paragraph
            if not has_content and not text_content:
                # Remove empty paragraph
                p._element.getparent().remove(p._element)
                continue

            # Apply paragraph alignment
            try:
                if paragraph.alignment is not None:
                    p.alignment = self._convert_alignment(paragraph.alignment)
            except:
                pass

            # Apply spacing
            try:
                if paragraph.space_before:
                    p.paragraph_format.space_before = paragraph.space_before
                if paragraph.space_after:
                    p.paragraph_format.space_after = paragraph.space_after
            except:
                pass

    def _convert_alignment(self, ppt_alignment):
        """Convert PowerPoint alignment to Word alignment"""
        alignment_map = {
            PP_PARAGRAPH_ALIGNMENT.LEFT: WD_ALIGN_PARAGRAPH.LEFT,
            PP_PARAGRAPH_ALIGNMENT.CENTER: WD_ALIGN_PARAGRAPH.CENTER,
            PP_PARAGRAPH_ALIGNMENT.RIGHT: WD_ALIGN_PARAGRAPH.RIGHT,
            PP_PARAGRAPH_ALIGNMENT.JUSTIFY: WD_ALIGN_PARAGRAPH.JUSTIFY,
        }
        return alignment_map.get(ppt_alignment, WD_ALIGN_PARAGRAPH.LEFT)

    def _extract_table_clean(self, shape, doc):
        """Extract table cleanly"""
        try:
            ppt_table = shape.table
            rows = len(ppt_table.rows)
            cols = len(ppt_table.columns)

            self._log(f"  Extracting table ({rows}×{cols})")

            # Create table in document
            doc_table = doc.add_table(rows=rows, cols=cols)
            doc_table.style = 'Light Grid Accent 1'

            # Set column widths if possible
            try:
                for i, col in enumerate(ppt_table.columns):
                    width_inches = col.width / 914400
                    doc_table.columns[i].width = Inches(width_inches)
            except:
                pass

            for i, row in enumerate(ppt_table.rows):
                for j, cell in enumerate(row.cells):
                    try:
                        doc_cell = doc_table.rows[i].cells[j]
                        doc_cell.text = ""

                        # Extract text with formatting
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                if len(doc_cell.paragraphs) == 0:
                                    cell_para = doc_cell.add_paragraph()
                                else:
                                    cell_para = doc_cell.paragraphs[0] if paragraph == cell.text_frame.paragraphs[0] else doc_cell.add_paragraph()

                                for run in paragraph.runs:
                                    cell_run = cell_para.add_run(run.text)

                                    # Copy formatting
                                    try:
                                        if run.font.bold:
                                            cell_run.font.bold = True
                                        if run.font.italic:
                                            cell_run.font.italic = True
                                        if run.font.size:
                                            cell_run.font.size = run.font.size
                                        if run.font.color and run.font.color.type == 1:
                                            rgb = run.font.color.rgb
                                            if rgb:
                                                cell_run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                                    except:
                                        pass

                                # Apply alignment
                                try:
                                    if paragraph.alignment:
                                        cell_para.alignment = self._convert_alignment(paragraph.alignment)
                                except:
                                    pass

                        # Copy cell background color if available
                        try:
                            if cell.fill.type == 1:  # Solid fill
                                rgb = cell.fill.fore_color.rgb
                                if rgb:
                                    # Set cell shading
                                    cell_xml = doc_cell._element
                                    cell_properties = cell_xml.get_or_add_tcPr()
                                    shd = OxmlElement('w:shd')
                                    shd.set(qn('w:fill'), '%02x%02x%02x' % (rgb[0], rgb[1], rgb[2]))
                                    cell_properties.append(shd)
                        except:
                            pass

                    except Exception as e:
                        self._log(f"  Warning: Could not process cell ({i},{j}): {e}")
                        continue

            # Add spacing after table
            doc.add_paragraph()

        except Exception as e:
            self._log(f"  Warning: Could not extract table: {e}")

    def _extract_picture_clean(self, shape, doc):
        """Extract picture cleanly"""
        try:
            image = shape.image
            image_bytes = image.blob

            # Preserve original dimensions
            width_inches = shape.width / 914400
            height_inches = shape.height / 914400

            # Only scale if too large
            max_width = 6.5
            if width_inches > max_width:
                ratio = max_width / width_inches
                width_inches = max_width
                height_inches = height_inches * ratio

            self._log(f"  Embedding image ({width_inches:.2f}\" × {height_inches:.2f}\")")

            # Center the image
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run()
            run.add_picture(io.BytesIO(image_bytes), width=Inches(width_inches))

            doc.add_paragraph()  # Spacing

        except Exception as e:
            self._log(f"  Warning: Could not extract image: {e}")

    def _extract_chart_clean(self, shape, doc):
        """Extract chart as data table if possible"""
        try:
            if not shape.has_chart:
                return

            chart = shape.chart
            chart_title = chart.chart_title.text_frame.text if chart.has_title else ""

            self._log(f"  Processing chart: {chart_title}")

            # Add chart title
            if chart_title:
                p = doc.add_paragraph()
                run = p.add_run(chart_title)
                run.font.bold = True
                run.font.size = Pt(14)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Try to extract chart data
            chart_data = self._extract_chart_data(chart)
            if chart_data and chart_data.get('categories') and chart_data.get('series'):
                # Create table for chart data
                rows = len(chart_data['categories']) + 1
                cols = len(chart_data['series']) + 1

                data_table = doc.add_table(rows=rows, cols=cols)
                data_table.style = 'Light List Accent 1'

                # Header row
                data_table.rows[0].cells[0].text = "Category"
                for i, series in enumerate(chart_data['series']):
                    cell = data_table.rows[0].cells[i + 1]
                    cell.text = series['name']
                    # Make header bold
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True

                # Data rows
                for row_idx, category in enumerate(chart_data['categories']):
                    data_table.rows[row_idx + 1].cells[0].text = str(category)
                    for col_idx, series in enumerate(chart_data['series']):
                        value = series['values'][row_idx] if row_idx < len(series['values']) else ""
                        data_table.rows[row_idx + 1].cells[col_idx + 1].text = str(value)

                doc.add_paragraph()
            else:
                # Add placeholder if can't extract data
                p = doc.add_paragraph()
                run = p.add_run("[Chart - data not available]")
                run.font.italic = True
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(128, 128, 128)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        except Exception as e:
            self._log(f"  Warning: Could not extract chart: {e}")

    def _extract_chart_data(self, chart):
        """Extract chart data"""
        try:
            chart_data = {
                'categories': [],
                'series': []
            }

            # Extract categories
            if hasattr(chart, 'plots') and len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, 'categories'):
                    chart_data['categories'] = [str(cat) for cat in plot.categories]

            # Extract series data
            for series in chart.series:
                series_data = {
                    'name': series.name or 'Series',
                    'values': []
                }
                for v in series.values:
                    if v is not None:
                        # Format numbers nicely
                        if isinstance(v, float):
                            series_data['values'].append(f"{v:.2f}")
                        else:
                            series_data['values'].append(str(v))
                    else:
                        series_data['values'].append("")

                chart_data['series'].append(series_data)

            return chart_data if chart_data['series'] else None

        except:
            return None
