from __future__ import annotations

import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.chart import XL_CHART_TYPE
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.shared import OxmlElement
from docx.oxml.ns import qn
from docx.oxml import parse_xml

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import numpy as np
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False


class ConversionError(Exception):
    pass


class PPTToDocConverterDitto:
    """
    Ditto converter - aims for maximum visual similarity to original presentation
    Uses absolute positioning and precise formatting
    """

    def __init__(self, progress_cb=None, log_cb=None):
        self._progress = progress_cb or (lambda p, m: None)
        self._log = log_cb or (lambda m: None)

    def convert(self, input_path: str | Path, output_path: str | Path):
        input_path = Path(input_path)
        output_path = Path(output_path)

        if not input_path.exists():
            raise ConversionError(f"Input file not found: {input_path}")

        try:
            self._log(f"Loading presentation: {input_path.name}")
            prs = Presentation(str(input_path))
        except Exception as e:
            raise ConversionError(f"Cannot open presentation: {e}") from e

        try:
            self._log("Creating Word document with ditto layout")
            doc = Document()

            # Match presentation dimensions exactly
            sections = doc.sections
            for section in sections:
                section.top_margin = Inches(0.3)
                section.bottom_margin = Inches(0.3)
                section.left_margin = Inches(0.3)
                section.right_margin = Inches(0.3)
                section.page_width = Inches(prs.slide_width / 914400)
                section.page_height = Inches(prs.slide_height / 914400)

            total_slides = len(prs.slides)
            self._log(f"Processing {total_slides} slides")

            for idx, slide in enumerate(prs.slides, 1):
                self._progress(int((idx / total_slides) * 90), f"Processing slide {idx}/{total_slides}")
                self._log(f"Processing slide {idx}")

                # Process entire slide as a layered layout
                self._process_slide_ditto(doc, slide, idx)

                # Page break
                if idx < total_slides:
                    doc.add_page_break()

            self._progress(95, "Saving document")
            self._log(f"Saving to: {output_path.name}")
            doc.save(str(output_path))

            self._progress(100, "Conversion complete")
            self._log("Conversion completed successfully")

        except Exception as e:
            raise ConversionError(f"Conversion failed: {e}") from e

    def _process_slide_ditto(self, doc, slide, slide_num):
        """Process slide maintaining visual layout"""

        # Create a container for this slide
        # Use a single-cell table as a positioning container
        container_table = doc.add_table(rows=1, cols=1)
        container_table.autofit = False
        container_table.allow_autofit = False

        # Set table width to match slide
        # Get width from presentation, not slide layout
        try:
            prs = slide.part.package.presentation_part.presentation
            slide_width_inches = prs.slide_width / 914400
        except:
            slide_width_inches = 10  # Default

        container_table.columns[0].width = Inches(slide_width_inches - 0.6)  # Account for margins

        cell = container_table.rows[0].cells[0]

        # Remove cell borders for clean look
        self._remove_cell_borders(cell)

        # Sort shapes by Z-order (background to foreground)
        sorted_shapes = self._sort_shapes_by_zorder(slide.shapes)

        # Process shapes in order
        for shape in sorted_shapes:
            try:
                # Skip embedded objects and connectors
                if shape.shape_type in [MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, MSO_SHAPE_TYPE.LINE]:
                    continue

                # Handle different shape types
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._process_group_ditto(cell, shape)
                elif hasattr(shape, 'text_frame') and shape.has_text_frame and shape.text_frame.text.strip():
                    self._add_text_shape_ditto(cell, shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    self._add_table_ditto(cell, shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._add_picture_ditto(cell, shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    self._add_chart_ditto(cell, shape)

            except Exception as e:
                self._log(f"  Warning: Could not process shape {shape.name}: {e}")
                continue

    def _sort_shapes_by_zorder(self, shapes):
        """Sort shapes by position (top to bottom, left to right)"""
        shape_list = []
        for shape in shapes:
            try:
                shape_list.append((shape.top, shape.left, shape))
            except:
                shape_list.append((float('inf'), float('inf'), shape))
        shape_list.sort(key=lambda x: (x[0], x[1]))
        return [s[2] for s in shape_list]

    def _process_group_ditto(self, cell, group_shape):
        """Process grouped shapes"""
        sorted_shapes = self._sort_shapes_by_zorder(group_shape.shapes)
        for shape in sorted_shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._process_group_ditto(cell, shape)
                elif hasattr(shape, 'text_frame') and shape.has_text_frame and shape.text_frame.text.strip():
                    self._add_text_shape_ditto(cell, shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    self._add_chart_ditto(cell, shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    self._add_table_ditto(cell, shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._add_picture_ditto(cell, shape)
            except Exception as e:
                continue

    def _add_text_shape_ditto(self, cell, shape):
        """Add text with exact formatting including background fills"""
        text_frame = shape.text_frame

        # Check if shape has a background fill (colored box)
        has_background_fill = False
        background_rgb = None
        try:
            if hasattr(shape, 'fill') and shape.fill.type == 1:  # Solid fill
                background_rgb = shape.fill.fore_color.rgb
                has_background_fill = True
        except:
            pass

        # If shape has background fill, use a single-cell table for the colored box
        if has_background_fill and background_rgb:
            # Create a single-cell table for the colored box
            box_table = cell.add_table(rows=1, cols=1)
            box_table.autofit = False
            box_table.allow_autofit = False

            box_cell = box_table.rows[0].cells[0]

            # Set cell background color
            hex_color = f'{background_rgb[0]:02x}{background_rgb[1]:02x}{background_rgb[2]:02x}'
            shading_elm = OxmlElement('w:shd')
            shading_elm.set(qn('w:fill'), hex_color)
            box_cell._element.get_or_add_tcPr().append(shading_elm)

            # Remove cell borders for clean look
            self._remove_cell_borders(box_cell)

            # Add text inside the colored box
            for paragraph in text_frame.paragraphs:
                if not paragraph.text.strip():
                    continue

                p = box_cell.add_paragraph()

                # Copy alignment
                try:
                    if paragraph.alignment == PP_PARAGRAPH_ALIGNMENT.CENTER:
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    elif paragraph.alignment == PP_PARAGRAPH_ALIGNMENT.RIGHT:
                        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                    elif paragraph.alignment == PP_PARAGRAPH_ALIGNMENT.JUSTIFY:
                        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    else:
                        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                except:
                    pass

                # Add runs with formatting
                self._add_runs_to_paragraph(p, paragraph)

            # Add spacing after the colored box
            cell.add_paragraph()
            return

        # Regular text without background fill
        for paragraph in text_frame.paragraphs:
            if not paragraph.text.strip():
                continue

            p = cell.add_paragraph()

            # Copy paragraph formatting
            try:
                if paragraph.alignment == PP_PARAGRAPH_ALIGNMENT.CENTER:
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                elif paragraph.alignment == PP_PARAGRAPH_ALIGNMENT.RIGHT:
                    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                elif paragraph.alignment == PP_PARAGRAPH_ALIGNMENT.JUSTIFY:
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                else:
                    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            except:
                pass

            # Add runs with formatting
            self._add_runs_to_paragraph(p, paragraph)

            # Space after paragraph
            p.paragraph_format.space_after = Pt(6)

    def _add_runs_to_paragraph(self, p, paragraph):
        """Add runs from PPT paragraph to Word paragraph with exact formatting"""
        for run in paragraph.runs:
            if not run.text:
                continue

            doc_run = p.add_run(run.text)

            # Font name - handle THEME/INHERIT
            if run.font.name:
                doc_run.font.name = run.font.name
            else:
                # If no explicit font, use a reasonable default
                # Titles/placeholders often use theme fonts
                doc_run.font.name = 'Calibri'  # Common theme default

            # Font size - preserve EXACT size or use reasonable default
            if run.font.size:
                doc_run.font.size = run.font.size
            else:
                # THEME/INHERIT size - use reasonable default
                # Placeholders/titles typically 14-18pt
                doc_run.font.size = Pt(14)

            # Bold/Italic/Underline - handle None (inherit) vs False (explicitly not)
            if run.font.bold is True:
                doc_run.font.bold = True
            elif run.font.bold is False:
                doc_run.font.bold = False
            # If None, leave as default

            if run.font.italic is True:
                doc_run.font.italic = True
            elif run.font.italic is False:
                doc_run.font.italic = False

            if run.font.underline is True:
                doc_run.font.underline = True
            elif run.font.underline is False:
                doc_run.font.underline = False

            # Color - handle RGB colors (including white!)
            if run.font.color and run.font.color.type == 1:  # RGB color type
                rgb = run.font.color.rgb
                if rgb:
                    # Set color even if it's white (255,255,255)
                    doc_run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
            # If color is THEME/INHERIT (type=None), use default black
            elif not run.font.color or run.font.color.type != 1:
                # Use default black for theme colors
                doc_run.font.color.rgb = RGBColor(0, 0, 0)

    def _add_table_ditto(self, cell, shape):
        """Add table with exact formatting"""
        ppt_table = shape.table
        rows = len(ppt_table.rows)
        cols = len(ppt_table.columns)

        self._log(f"  Adding table ({rows}×{cols})")

        # Add table to cell
        doc_table = cell.add_table(rows=rows, cols=cols)
        doc_table.style = 'Light Grid'

        # Copy cell content and formatting
        for i, row in enumerate(ppt_table.rows):
            for j, ppt_cell in enumerate(row.cells):
                doc_cell = doc_table.rows[i].cells[j]

                # Clear default content
                doc_cell.text = ""

                # Copy text with formatting
                if ppt_cell.text_frame:
                    for para in ppt_cell.text_frame.paragraphs:
                        cell_para = doc_cell.paragraphs[0] if len(doc_cell.paragraphs) > 0 else doc_cell.add_paragraph()

                        # Copy alignment
                        try:
                            if para.alignment == PP_PARAGRAPH_ALIGNMENT.CENTER:
                                cell_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            elif para.alignment == PP_PARAGRAPH_ALIGNMENT.RIGHT:
                                cell_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                        except:
                            pass

                        # Copy runs
                        for run in para.runs:
                            cell_run = cell_para.add_run(run.text)

                            if run.font.name:
                                cell_run.font.name = run.font.name
                            if run.font.size:
                                cell_run.font.size = run.font.size
                            if run.font.bold:
                                cell_run.font.bold = True
                            if run.font.italic:
                                cell_run.font.italic = True
                            if run.font.color and run.font.color.type == 1:
                                rgb = run.font.color.rgb
                                if rgb:
                                    cell_run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])

                # Copy cell background
                try:
                    if ppt_cell.fill.type == 1:  # Solid
                        rgb = ppt_cell.fill.fore_color.rgb
                        if rgb:
                            shading_elm = parse_xml(f'<w:shd {qn("w:fill")}="{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"/>')
                            doc_cell._element.get_or_add_tcPr().append(shading_elm)
                except:
                    pass

        # Add spacing
        cell.add_paragraph()

    def _add_picture_ditto(self, cell, shape):
        """Add picture preserving size"""
        try:
            image_bytes = shape.image.blob
            width_inches = shape.width / 914400
            height_inches = shape.height / 914400

            # Cap at reasonable size
            max_width = 6.5
            if width_inches > max_width:
                ratio = max_width / width_inches
                width_inches = max_width
                height_inches = height_inches * ratio

            self._log(f"  Adding image ({width_inches:.2f}\" × {height_inches:.2f}\")")

            p = cell.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run()
            run.add_picture(io.BytesIO(image_bytes), width=Inches(width_inches))

            cell.add_paragraph()

        except Exception as e:
            self._log(f"  Warning: Could not add image: {e}")

    def _add_chart_ditto(self, cell, shape):
        """Add chart as visual image"""
        if not MATPLOTLIB_AVAILABLE:
            self._log("  Warning: matplotlib not available")
            return

        try:
            if not shape.has_chart:
                return

            chart = shape.chart
            chart_type = chart.chart_type

            self._log(f"  Rendering chart (type: {chart_type})")

            # Extract data
            chart_data = self._extract_chart_data(chart)
            if not chart_data or not chart_data.get('series'):
                return

            categories = chart_data.get('categories', [])
            series_list = chart_data['series']

            # Create figure with appropriate size
            width_inches = shape.width / 914400
            height_inches = shape.height / 914400

            # Cap size
            max_width = 7
            if width_inches > max_width:
                ratio = max_width / width_inches
                width_inches = max_width
                height_inches = height_inches * ratio

            fig, ax = plt.subplots(figsize=(width_inches, height_inches), dpi=120)

            # Render chart
            if chart_type in [XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.BAR_CLUSTERED]:
                self._render_bar_chart(ax, categories, series_list, chart_type)
            elif chart_type in [XL_CHART_TYPE.COLUMN_STACKED, XL_CHART_TYPE.COLUMN_STACKED_100]:
                self._render_stacked_bar_chart(ax, categories, series_list, chart_type)
            elif chart_type in [XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS, XL_CHART_TYPE.LINE_MARKERS_STACKED]:
                self._render_line_chart(ax, categories, series_list)
            elif chart_type in [XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED]:
                self._render_pie_chart(ax, categories, series_list)
            else:
                self._render_bar_chart(ax, categories, series_list, chart_type)

            # Title
            if chart.has_title:
                try:
                    title = chart.chart_title.text_frame.text
                    if title:
                        ax.set_title(title, fontsize=12, fontweight='bold', pad=10)
                except:
                    pass

            # Tight layout
            plt.tight_layout()

            # Save to bytes
            img_bytes = io.BytesIO()
            plt.savefig(img_bytes, format='png', bbox_inches='tight', dpi=150, facecolor='white')
            plt.close(fig)
            img_bytes.seek(0)

            # Add to document
            p = cell.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run()
            run.add_picture(img_bytes, width=Inches(width_inches))

            cell.add_paragraph()

        except Exception as e:
            self._log(f"  Warning: Could not render chart: {e}")

    def _extract_chart_data(self, chart):
        """Extract chart data"""
        try:
            chart_data = {'categories': [], 'series': []}

            if hasattr(chart, 'plots') and len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, 'categories'):
                    chart_data['categories'] = [str(cat) for cat in plot.categories]

            for series in chart.series:
                series_data = {
                    'name': series.name or 'Series',
                    'values': [v if v is not None else 0 for v in series.values]
                }
                chart_data['series'].append(series_data)

            return chart_data if chart_data['series'] else None
        except:
            return None

    def _render_bar_chart(self, ax, categories, series_list, chart_type):
        """Render bar chart"""
        x = np.arange(len(categories))
        width = 0.8 / max(len(series_list), 1)

        for i, series in enumerate(series_list):
            offset = (i - len(series_list)/2) * width + width/2
            ax.bar(x + offset, series['values'], width, label=series['name'], alpha=0.8)

        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=0 if len(categories) <= 4 else 45, ha='right' if len(categories) > 4 else 'center', fontsize=9)
        ax.legend(fontsize=8, loc='upper left')
        ax.grid(axis='y', alpha=0.3, linestyle='--')
        ax.set_axisbelow(True)

    def _render_stacked_bar_chart(self, ax, categories, series_list, chart_type):
        """Render stacked bar chart"""
        x = np.arange(len(categories))
        width = 0.7

        bottom = np.zeros(len(categories))

        for series in series_list:
            values = np.array(series['values'])
            ax.bar(x, values, width, label=series['name'], bottom=bottom, alpha=0.9)
            bottom += values

        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=0 if len(categories) <= 4 else 45, ha='right' if len(categories) > 4 else 'center', fontsize=9)
        ax.legend(fontsize=8, loc='upper left')
        ax.grid(axis='y', alpha=0.3, linestyle='--')
        ax.set_axisbelow(True)

        if chart_type == XL_CHART_TYPE.COLUMN_STACKED_100:
            ax.set_ylabel('Percentage (%)', fontsize=9)
            ax.set_ylim(0, 100)

    def _render_line_chart(self, ax, categories, series_list):
        """Render line chart"""
        x = np.arange(len(categories))

        for series in series_list:
            ax.plot(x, series['values'], marker='o', label=series['name'], linewidth=2)

        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right', fontsize=9)
        ax.legend(fontsize=8)
        ax.grid(alpha=0.3, linestyle='--')

    def _render_pie_chart(self, ax, categories, series_list):
        """Render pie chart"""
        if len(series_list) > 0:
            series = series_list[0]
            wedges, texts, autotexts = ax.pie(series['values'], labels=categories, autopct='%1.1f%%', startangle=90)
            for text in texts:
                text.set_fontsize(9)
            for autotext in autotexts:
                autotext.set_fontsize(8)
                autotext.set_color('white')
                autotext.set_weight('bold')
            ax.axis('equal')

    def _remove_cell_borders(self, cell):
        """Remove all borders from a table cell"""
        tc = cell._element
        tcPr = tc.get_or_add_tcPr()
        tcBorders = OxmlElement('w:tcBorders')
        for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
            border = OxmlElement(f'w:{border_name}')
            border.set(qn('w:val'), 'none')
            border.set(qn('w:sz'), '0')
            border.set(qn('w:space'), '0')
            border.set(qn('w:color'), 'auto')
            tcBorders.append(border)
        tcPr.append(tcBorders)
