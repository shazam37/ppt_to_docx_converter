from __future__ import annotations

import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml.ns import qn
from docx.oxml import OxmlElement


class ConversionError(Exception):
    pass


class PPTToDocConverterEnhanced:
    def __init__(self, progress_cb=None, log_cb=None):
        """
        Enhanced converter with maximum fidelity preservation
        progress_cb(percent: int, message: str)
        log_cb(message: str)
        """
        self._progress = progress_cb or (lambda p, m: None)
        self._log = log_cb or (lambda m: None)

    def convert(self, input_path: str | Path, output_path: str | Path):
        input_path = Path(input_path)
        output_path = Path(output_path)

        if not input_path.exists():
            raise ConversionError(f"Input file not found: {input_path}")

        try:
            self._log(f"Loading presentation: {input_path.name}")
            prs = Presentation(str(input_path))
        except Exception as e:
            raise ConversionError(f"Cannot open presentation: {e}") from e

        try:
            self._log("Creating Word document with enhanced fidelity")
            doc = Document()

            # Set up document margins (minimal for maximum space)
            sections = doc.sections
            for section in sections:
                section.top_margin = Inches(0.5)
                section.bottom_margin = Inches(0.5)
                section.left_margin = Inches(0.5)
                section.right_margin = Inches(0.5)
                # Set page size to match presentation aspect ratio
                section.page_width = Inches(10)
                section.page_height = Inches(7.5)

            total_slides = len(prs.slides)
            self._log(f"Processing {total_slides} slides")

            for idx, slide in enumerate(prs.slides, 1):
                self._progress(int((idx / total_slides) * 90), f"Processing slide {idx}/{total_slides}")
                self._log(f"Processing slide {idx}")

                # Add slide header with metadata
                self._add_slide_header(doc, idx, slide)

                # Process slide background if available
                self._process_slide_background(doc, slide)

                # Sort shapes by position (top-to-bottom, left-to-right) for natural reading order
                sorted_shapes = self._sort_shapes_by_position(slide.shapes)

                # Process all shapes in sorted order
                self._process_shapes(sorted_shapes, doc, slide)

                # Add slide footer with metadata
                self._add_slide_footer(doc, slide)

                # Add page break after each slide except the last one
                if idx < total_slides:
                    doc.add_page_break()

            self._progress(95, "Saving document")
            self._log(f"Saving to: {output_path.name}")
            doc.save(str(output_path))

            self._progress(100, "Conversion complete")
            self._log("Conversion completed successfully")

        except Exception as e:
            raise ConversionError(f"Conversion failed: {e}") from e

    def _add_slide_header(self, doc, slide_num, slide):
        """Add comprehensive slide header with metadata"""
        # Slide number
        heading = doc.add_heading(f"Slide {slide_num}", level=1)
        heading.alignment = WD_ALIGN_PARAGRAPH.LEFT

        # Layout information
        try:
            layout_name = slide.slide_layout.name
            if layout_name:
                p = doc.add_paragraph()
                run = p.add_run(f"📐 Layout: {layout_name}")
                run.font.size = Pt(9)
                run.font.italic = True
                run.font.color.rgb = RGBColor(100, 100, 100)
        except:
            pass

        # Slide dimensions
        try:
            width_inches = slide.slide_layout.slide_width / 914400
            height_inches = slide.slide_layout.slide_height / 914400
            p = doc.add_paragraph()
            run = p.add_run(f"📏 Dimensions: {width_inches:.2f}\" × {height_inches:.2f}\"")
            run.font.size = Pt(9)
            run.font.italic = True
            run.font.color.rgb = RGBColor(100, 100, 100)
        except:
            pass

        # Number of shapes
        p = doc.add_paragraph()
        run = p.add_run(f"🔢 Elements: {len(slide.shapes)} shape(s)")
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.rgb = RGBColor(100, 100, 100)

        # Separator
        doc.add_paragraph("─" * 80)

    def _add_slide_footer(self, doc, slide):
        """Add slide footer with notes if present"""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    doc.add_paragraph("─" * 80)
                    p = doc.add_paragraph()
                    run = p.add_run("📝 Speaker Notes:")
                    run.font.bold = True
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(50, 50, 150)

                    notes_p = doc.add_paragraph(notes_text)
                    notes_p.paragraph_format.left_indent = Inches(0.25)
                    for run in notes_p.runs:
                        run.font.size = Pt(9)
                        run.font.italic = True
                        run.font.color.rgb = RGBColor(80, 80, 80)
        except:
            pass

    def _process_slide_background(self, doc, slide):
        """Document slide background information"""
        try:
            background = slide.background
            if background.fill.type:
                p = doc.add_paragraph()
                run = p.add_run(f"🎨 Background: {background.fill.type}")
                run.font.size = Pt(9)
                run.font.italic = True
                run.font.color.rgb = RGBColor(100, 100, 100)
        except:
            pass

    def _sort_shapes_by_position(self, shapes):
        """Sort shapes by vertical then horizontal position for natural reading order"""
        shape_list = []
        for shape in shapes:
            try:
                # Get shape position (top, left)
                top = shape.top
                left = shape.left
                shape_list.append((top, left, shape))
            except:
                # If position unavailable, add to end
                shape_list.append((float('inf'), float('inf'), shape))

        # Sort by top first, then left
        shape_list.sort(key=lambda x: (x[0], x[1]))

        # Return just the shapes
        return [s[2] for s in shape_list]

    def _process_shapes(self, shapes, doc, slide, indent_level=0):
        """Process shapes with maximum detail preservation"""
        for shape in shapes:
            try:
                # Add shape metadata for complex layouts
                if indent_level == 0:  # Only for top-level shapes
                    self._add_shape_metadata(doc, shape)

                # Handle grouped shapes recursively
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._log(f"  Processing group with {len(shape.shapes)} shapes")
                    # Add group marker
                    p = doc.add_paragraph()
                    run = p.add_run(f"┌─ Group: {len(shape.shapes)} elements ─┐")
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(150, 150, 150)

                    sorted_group_shapes = self._sort_shapes_by_position(shape.shapes)
                    self._process_shapes(sorted_group_shapes, doc, slide, indent_level + 1)

                    p = doc.add_paragraph()
                    run = p.add_run("└─────────────────────┘")
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(150, 150, 150)
                    continue

                # Handle text boxes and shapes with text
                if shape.has_text_frame:
                    self._extract_text_frame_enhanced(shape, doc, indent_level)

                # Handle tables
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    self._extract_table_enhanced(shape, doc)

                # Handle pictures
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._extract_picture_enhanced(shape, doc)

                # Handle charts
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    self._extract_chart_enhanced(shape, doc)

                # Handle other shapes
                else:
                    self._handle_other_shape(shape, doc)

            except Exception as e:
                self._log(f"  Warning: Could not process shape: {e}")
                # Add error marker
                p = doc.add_paragraph()
                run = p.add_run(f"[⚠️ Shape processing error: {str(e)[:100]}]")
                run.font.size = Pt(8)
                run.font.color.rgb = RGBColor(200, 0, 0)
                continue

    def _add_shape_metadata(self, doc, shape):
        """Add metadata about shape for complex layout documentation"""
        try:
            # Position and size info
            left = shape.left / 914400
            top = shape.top / 914400
            width = shape.width / 914400
            height = shape.height / 914400

            p = doc.add_paragraph()
            run = p.add_run(f"📍 Position: ({left:.2f}\", {top:.2f}\") | Size: {width:.2f}\" × {height:.2f}\"")
            run.font.size = Pt(7)
            run.font.color.rgb = RGBColor(180, 180, 180)
        except:
            pass

    def _extract_text_frame_enhanced(self, shape, doc, indent_level=0):
        """Extract text with maximum formatting preservation"""
        text_frame = shape.text_frame

        # Add shape name if available
        try:
            if shape.name and not shape.name.startswith("Text"):
                p = doc.add_paragraph()
                run = p.add_run(f"🏷️ {shape.name}")
                run.font.size = Pt(8)
                run.font.italic = True
                run.font.color.rgb = RGBColor(120, 120, 120)
        except:
            pass

        for paragraph in text_frame.paragraphs:
            # Process even empty paragraphs to preserve spacing
            p = doc.add_paragraph()

            # Apply indentation
            if indent_level > 0:
                p.paragraph_format.left_indent = Inches(0.25 * indent_level)

            # Apply paragraph-level indentation from PPT
            try:
                if paragraph.level > 0:
                    p.paragraph_format.left_indent = Inches(0.25 * (indent_level + paragraph.level))
            except:
                pass

            # Process runs to preserve formatting
            has_content = False
            for run in paragraph.runs:
                text = run.text
                if not text:
                    continue

                has_content = True
                doc_run = p.add_run(text)

                # Copy all available formatting
                try:
                    # Font properties
                    if run.font.bold is not None:
                        doc_run.font.bold = run.font.bold
                    if run.font.italic is not None:
                        doc_run.font.italic = run.font.italic
                    if run.font.underline is not None:
                        doc_run.font.underline = run.font.underline

                    # Font size - preserve actual size
                    if run.font.size:
                        doc_run.font.size = run.font.size

                    # Font name
                    if run.font.name:
                        doc_run.font.name = run.font.name

                    # Text color
                    if run.font.color and run.font.color.type == 1:
                        rgb = run.font.color.rgb
                        if rgb:
                            doc_run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])

                    # Superscript/Subscript
                    if hasattr(run.font, 'superscript') and run.font.superscript:
                        doc_run.font.superscript = True
                    if hasattr(run.font, 'subscript') and run.font.subscript:
                        doc_run.font.subscript = True

                except Exception as e:
                    self._log(f"  Warning: Could not copy all formatting: {e}")

            # If paragraph has no content but exists, add placeholder
            if not has_content and paragraph.text.strip():
                doc_run = p.add_run(paragraph.text)

            # Apply paragraph alignment
            try:
                if paragraph.alignment is not None:
                    p.alignment = self._convert_alignment(paragraph.alignment)
            except:
                pass

            # Apply line spacing if available
            try:
                if paragraph.line_spacing:
                    p.paragraph_format.line_spacing = paragraph.line_spacing
            except:
                pass

            # Apply space before/after
            try:
                if paragraph.space_before:
                    p.paragraph_format.space_before = paragraph.space_before
                if paragraph.space_after:
                    p.paragraph_format.space_after = paragraph.space_after
            except:
                pass

    def _convert_alignment(self, ppt_alignment):
        """Convert PowerPoint alignment to Word alignment"""
        alignment_map = {
            PP_PARAGRAPH_ALIGNMENT.LEFT: WD_ALIGN_PARAGRAPH.LEFT,
            PP_PARAGRAPH_ALIGNMENT.CENTER: WD_ALIGN_PARAGRAPH.CENTER,
            PP_PARAGRAPH_ALIGNMENT.RIGHT: WD_ALIGN_PARAGRAPH.RIGHT,
            PP_PARAGRAPH_ALIGNMENT.JUSTIFY: WD_ALIGN_PARAGRAPH.JUSTIFY,
        }
        return alignment_map.get(ppt_alignment, WD_ALIGN_PARAGRAPH.LEFT)

    def _extract_table_enhanced(self, shape, doc):
        """Extract table with enhanced formatting preservation"""
        try:
            ppt_table = shape.table
            rows = len(ppt_table.rows)
            cols = len(ppt_table.columns)

            self._log(f"  Extracting table ({rows}×{cols})")

            # Add table marker
            p = doc.add_paragraph()
            run = p.add_run(f"📊 Table: {rows} rows × {cols} columns")
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 100, 0)

            # Create table in document
            doc_table = doc.add_table(rows=rows, cols=cols)
            doc_table.style = 'Light Grid Accent 1'

            for i, row in enumerate(ppt_table.rows):
                for j, cell in enumerate(row.cells):
                    try:
                        # Extract cell text with formatting
                        cell_text = cell.text.strip()
                        doc_cell = doc_table.rows[i].cells[j]

                        # Clear default text
                        doc_cell.text = ""

                        # Add text with formatting from text frame
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                cell_para = doc_cell.paragraphs[0] if len(doc_cell.paragraphs) > 0 else doc_cell.add_paragraph()

                                for run in paragraph.runs:
                                    cell_run = cell_para.add_run(run.text)

                                    # Copy formatting
                                    try:
                                        if run.font.bold:
                                            cell_run.font.bold = True
                                        if run.font.italic:
                                            cell_run.font.italic = True
                                        if run.font.size:
                                            cell_run.font.size = run.font.size
                                        if run.font.color and run.font.color.type == 1:
                                            rgb = run.font.color.rgb
                                            if rgb:
                                                cell_run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                                    except:
                                        pass

                        # Special formatting for header row
                        if i == 0:
                            for paragraph in doc_cell.paragraphs:
                                for run in paragraph.runs:
                                    run.font.bold = True

                    except Exception as e:
                        self._log(f"  Warning: Could not process cell ({i},{j}): {e}")
                        continue

            # Add spacing after table
            doc.add_paragraph()

        except Exception as e:
            self._log(f"  Warning: Could not extract table: {e}")

    def _extract_picture_enhanced(self, shape, doc):
        """Extract and embed picture with metadata"""
        try:
            image = shape.image
            image_bytes = image.blob

            # Get image properties
            width = shape.width / 914400
            height = shape.height / 914400
            content_type = image.content_type

            self._log(f"  Embedding image: {content_type} ({width:.2f}\" × {height:.2f}\")")

            # Add image marker
            p = doc.add_paragraph()
            run = p.add_run(f"🖼️ Image: {content_type} | {width:.2f}\" × {height:.2f}\"")
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(0, 100, 200)

            # Preserve original size if reasonable, otherwise scale
            max_width = 6.5
            if width > max_width:
                ratio = max_width / width
                width = max_width
                height = height * ratio

            # Add image to document
            doc.add_picture(io.BytesIO(image_bytes), width=Inches(width))

            # Add shape name if available
            try:
                if shape.name and not shape.name.startswith("Picture"):
                    p = doc.add_paragraph()
                    run = p.add_run(f"Caption: {shape.name}")
                    run.font.size = Pt(8)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(100, 100, 100)
            except:
                pass

            doc.add_paragraph()  # Add spacing

        except Exception as e:
            self._log(f"  Warning: Could not extract image: {e}")
            # Add placeholder with error
            p = doc.add_paragraph()
            run = p.add_run(f"[❌ Image extraction failed: {str(e)[:100]}]")
            run.font.italic = True
            run.font.color.rgb = RGBColor(255, 0, 0)

    def _extract_chart_enhanced(self, shape, doc):
        """Extract chart with data extraction attempt"""
        try:
            if shape.has_chart:
                chart = shape.chart
                chart_title = chart.chart_title.text_frame.text if chart.has_title else "Untitled Chart"

                self._log(f"  Processing chart: {chart_title}")

                # Chart header
                p = doc.add_paragraph()
                run = p.add_run(f"📈 Chart: {chart_title}")
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0, 0, 200)

                # Try to extract chart data
                try:
                    chart_data = self._extract_chart_data(chart)
                    if chart_data:
                        # Add chart data as table
                        p = doc.add_paragraph()
                        run = p.add_run("Chart Data:")
                        run.font.bold = True
                        run.font.size = Pt(10)

                        # Create table for chart data
                        if chart_data.get('categories') and chart_data.get('series'):
                            rows = len(chart_data['series'][0]['values']) + 1
                            cols = len(chart_data['series']) + 1

                            data_table = doc.add_table(rows=rows, cols=cols)
                            data_table.style = 'Light List Accent 1'

                            # Header row
                            data_table.rows[0].cells[0].text = "Category"
                            for i, series in enumerate(chart_data['series']):
                                data_table.rows[0].cells[i + 1].text = series['name']

                            # Data rows
                            for row_idx, category in enumerate(chart_data['categories']):
                                data_table.rows[row_idx + 1].cells[0].text = str(category)
                                for col_idx, series in enumerate(chart_data['series']):
                                    value = series['values'][row_idx]
                                    data_table.rows[row_idx + 1].cells[col_idx + 1].text = str(value)

                            doc.add_paragraph()
                except Exception as e:
                    self._log(f"  Could not extract chart data: {e}")
                    p = doc.add_paragraph()
                    run = p.add_run("(Chart data extraction not available - refer to original presentation)")
                    run.font.size = Pt(9)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(150, 150, 150)

        except Exception as e:
            self._log(f"  Warning: Could not extract chart: {e}")

    def _extract_chart_data(self, chart):
        """Attempt to extract chart data"""
        try:
            chart_data = {
                'categories': [],
                'series': []
            }

            # Extract categories
            if hasattr(chart, 'plots') and len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, 'categories'):
                    chart_data['categories'] = [str(cat) for cat in plot.categories]

            # Extract series data
            for series in chart.series:
                series_data = {
                    'name': series.name or 'Series',
                    'values': [v if v is not None else 0 for v in series.values]
                }
                chart_data['series'].append(series_data)

            return chart_data if chart_data['series'] else None

        except:
            return None

    def _handle_other_shape(self, shape, doc):
        """Handle other shape types with documentation"""
        try:
            shape_type = shape.shape_type
            shape_name = shape.name if hasattr(shape, 'name') else "Unknown"

            self._log(f"  Found shape: {shape_name} (type: {shape_type})")

            p = doc.add_paragraph()
            run = p.add_run(f"[🔷 Shape: {shape_name} | Type: {shape_type}]")
            run.font.size = Pt(9)
            run.font.italic = True
            run.font.color.rgb = RGBColor(100, 100, 200)

        except Exception as e:
            self._log(f"  Warning: Could not document shape: {e}")
