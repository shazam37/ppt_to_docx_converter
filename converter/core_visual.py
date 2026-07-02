from __future__ import annotations

import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.chart import XL_CHART_TYPE
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.shared import OxmlElement
from docx.oxml.ns import qn

try:
    import matplotlib
    matplotlib.use('Agg')  # Non-interactive backend
    import matplotlib.pyplot as plt
    import numpy as np
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False


class ConversionError(Exception):
    pass


class PPTToDocConverterVisual:
    def __init__(self, progress_cb=None, log_cb=None):
        """
        Visual converter - preserves charts as visual images
        progress_cb(percent: int, message: str)
        log_cb(message: str)
        """
        self._progress = progress_cb or (lambda p, m: None)
        self._log = log_cb or (lambda m: None)

    def convert(self, input_path: str | Path, output_path: str | Path):
        input_path = Path(input_path)
        output_path = Path(output_path)

        if not input_path.exists():
            raise ConversionError(f"Input file not found: {input_path}")

        try:
            self._log(f"Loading presentation: {input_path.name}")
            prs = Presentation(str(input_path))
        except Exception as e:
            raise ConversionError(f"Cannot open presentation: {e}") from e

        try:
            self._log("Creating Word document with visual charts")
            doc = Document()

            # Set up document to match presentation aspect ratio
            sections = doc.sections
            for section in sections:
                section.top_margin = Inches(0.5)
                section.bottom_margin = Inches(0.5)
                section.left_margin = Inches(0.75)
                section.right_margin = Inches(0.75)
                section.page_width = Inches(prs.slide_width / 914400)
                section.page_height = Inches(prs.slide_height / 914400)

            total_slides = len(prs.slides)
            self._log(f"Processing {total_slides} slides")

            for idx, slide in enumerate(prs.slides, 1):
                self._progress(int((idx / total_slides) * 90), f"Processing slide {idx}/{total_slides}")
                self._log(f"Processing slide {idx}")

                # Slide header
                heading = doc.add_heading(f"Slide {idx}", level=1)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

                # Separator
                p = doc.add_paragraph("─" * 80)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER

                # Sort shapes by position
                sorted_shapes = self._sort_shapes_by_position(slide.shapes)

                # Process shapes
                self._process_shapes(sorted_shapes, doc, slide)

                # Page break
                if idx < total_slides:
                    doc.add_page_break()

            self._progress(95, "Saving document")
            self._log(f"Saving to: {output_path.name}")
            doc.save(str(output_path))

            self._progress(100, "Conversion complete")
            self._log("Conversion completed successfully")

        except Exception as e:
            raise ConversionError(f"Conversion failed: {e}") from e

    def _sort_shapes_by_position(self, shapes):
        """Sort shapes by position"""
        shape_list = []
        for shape in shapes:
            try:
                shape_list.append((shape.top, shape.left, shape))
            except:
                shape_list.append((float('inf'), float('inf'), shape))
        shape_list.sort(key=lambda x: (x[0], x[1]))
        return [s[2] for s in shape_list]

    def _process_shapes(self, shapes, doc, slide, indent_level=0):
        """Process shapes"""
        for shape in shapes:
            try:
                # Skip embedded objects and lines
                if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    continue

                # Handle grouped shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._log(f"  Processing group with {len(shape.shapes)} shapes")
                    sorted_group = self._sort_shapes_by_position(shape.shapes)
                    self._process_shapes(sorted_group, doc, slide, indent_level)
                    continue

                # Handle text
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    text_content = shape.text_frame.text.strip()
                    if text_content:
                        self._extract_text_frame(shape, doc, indent_level)

                # Handle tables
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    self._extract_table(shape, doc)

                # Handle pictures
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._extract_picture(shape, doc)

                # Handle charts - AS VISUAL IMAGES
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    self._extract_chart_visual(shape, doc)

            except Exception as e:
                self._log(f"  Warning: Could not process shape: {e}")
                continue

    def _extract_text_frame(self, shape, doc, indent_level=0):
        """Extract text"""
        text_frame = shape.text_frame

        for paragraph in text_frame.paragraphs:
            text_content = paragraph.text.strip()

            p = doc.add_paragraph()

            # Indentation
            total_indent = indent_level
            try:
                if paragraph.level > 0:
                    total_indent += paragraph.level
            except:
                pass

            if total_indent > 0:
                p.paragraph_format.left_indent = Inches(0.25 * total_indent)

            # Process runs
            has_content = False
            for run in paragraph.runs:
                if not run.text:
                    continue

                has_content = True
                doc_run = p.add_run(run.text)

                # Formatting
                try:
                    if run.font.bold is not None:
                        doc_run.font.bold = run.font.bold
                    if run.font.italic is not None:
                        doc_run.font.italic = run.font.italic
                    if run.font.underline is not None:
                        doc_run.font.underline = run.font.underline
                    if run.font.size:
                        doc_run.font.size = run.font.size
                    if run.font.name:
                        doc_run.font.name = run.font.name
                    if run.font.color and run.font.color.type == 1:
                        rgb = run.font.color.rgb
                        if rgb:
                            doc_run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                except:
                    pass

            # Remove empty paragraphs
            if not has_content and not text_content:
                p._element.getparent().remove(p._element)
                continue

            # Alignment
            try:
                if paragraph.alignment is not None:
                    p.alignment = self._convert_alignment(paragraph.alignment)
            except:
                pass

    def _convert_alignment(self, ppt_alignment):
        """Convert alignment"""
        alignment_map = {
            PP_PARAGRAPH_ALIGNMENT.LEFT: WD_ALIGN_PARAGRAPH.LEFT,
            PP_PARAGRAPH_ALIGNMENT.CENTER: WD_ALIGN_PARAGRAPH.CENTER,
            PP_PARAGRAPH_ALIGNMENT.RIGHT: WD_ALIGN_PARAGRAPH.RIGHT,
            PP_PARAGRAPH_ALIGNMENT.JUSTIFY: WD_ALIGN_PARAGRAPH.JUSTIFY,
        }
        return alignment_map.get(ppt_alignment, WD_ALIGN_PARAGRAPH.LEFT)

    def _extract_table(self, shape, doc):
        """Extract table"""
        try:
            ppt_table = shape.table
            rows = len(ppt_table.rows)
            cols = len(ppt_table.columns)

            self._log(f"  Extracting table ({rows}×{cols})")

            doc_table = doc.add_table(rows=rows, cols=cols)
            doc_table.style = 'Light Grid Accent 1'

            for i, row in enumerate(ppt_table.rows):
                for j, cell in enumerate(row.cells):
                    try:
                        doc_cell = doc_table.rows[i].cells[j]
                        doc_cell.text = ""

                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                cell_para = doc_cell.paragraphs[0] if len(doc_cell.paragraphs) > 0 else doc_cell.add_paragraph()

                                for run in paragraph.runs:
                                    cell_run = cell_para.add_run(run.text)

                                    try:
                                        if run.font.bold:
                                            cell_run.font.bold = True
                                        if run.font.italic:
                                            cell_run.font.italic = True
                                        if run.font.size:
                                            cell_run.font.size = run.font.size
                                        if run.font.color and run.font.color.type == 1:
                                            rgb = run.font.color.rgb
                                            if rgb:
                                                cell_run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                                    except:
                                        pass

                    except Exception as e:
                        continue

            doc.add_paragraph()

        except Exception as e:
            self._log(f"  Warning: Could not extract table: {e}")

    def _extract_picture(self, shape, doc):
        """Extract picture"""
        try:
            image = shape.image
            image_bytes = image.blob

            width_inches = shape.width / 914400
            height_inches = shape.height / 914400

            max_width = 6.5
            if width_inches > max_width:
                ratio = max_width / width_inches
                width_inches = max_width
                height_inches = height_inches * ratio

            self._log(f"  Embedding image ({width_inches:.2f}\" × {height_inches:.2f}\")")

            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run()
            run.add_picture(io.BytesIO(image_bytes), width=Inches(width_inches))

            doc.add_paragraph()

        except Exception as e:
            self._log(f"  Warning: Could not extract image: {e}")

    def _extract_chart_visual(self, shape, doc):
        """Extract chart as VISUAL IMAGE using matplotlib"""
        try:
            if not shape.has_chart:
                return

            if not MATPLOTLIB_AVAILABLE:
                self._log("  Warning: matplotlib not available, adding placeholder")
                p = doc.add_paragraph()
                run = p.add_run("[Chart - matplotlib required for visualization]")
                run.font.italic = True
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(200, 0, 0)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                return

            chart = shape.chart
            chart_type = chart.chart_type

            self._log(f"  Rendering chart as image (type: {chart_type})")

            # Create matplotlib figure
            fig, ax = plt.subplots(figsize=(8, 5), dpi=100)

            # Extract chart data
            chart_data = self._extract_chart_data(chart)

            if not chart_data or not chart_data.get('series'):
                self._log("  Warning: Could not extract chart data")
                plt.close(fig)
                return

            categories = chart_data.get('categories', [])
            series_list = chart_data['series']

            # Render based on chart type
            if chart_type in [XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.BAR_CLUSTERED]:
                self._render_bar_chart(ax, categories, series_list, chart_type)
            elif chart_type in [XL_CHART_TYPE.COLUMN_STACKED, XL_CHART_TYPE.COLUMN_STACKED_100]:
                self._render_stacked_bar_chart(ax, categories, series_list, chart_type)
            elif chart_type in [XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS, XL_CHART_TYPE.LINE_MARKERS_STACKED]:
                self._render_line_chart(ax, categories, series_list)
            elif chart_type in [XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED]:
                self._render_pie_chart(ax, categories, series_list)
            else:
                # Default to bar chart
                self._render_bar_chart(ax, categories, series_list, chart_type)

            # Add title if available
            if chart.has_title:
                try:
                    title = chart.chart_title.text_frame.text
                    if title:
                        ax.set_title(title, fontsize=14, fontweight='bold')
                except:
                    pass

            plt.tight_layout()

            # Save to bytes
            img_bytes = io.BytesIO()
            plt.savefig(img_bytes, format='png', bbox_inches='tight', dpi=150)
            plt.close(fig)
            img_bytes.seek(0)

            # Add to document
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run()
            run.add_picture(img_bytes, width=Inches(6.5))

            doc.add_paragraph()

        except Exception as e:
            self._log(f"  Warning: Could not render chart: {e}")
            import traceback
            traceback.print_exc()

    def _extract_chart_data(self, chart):
        """Extract chart data"""
        try:
            chart_data = {'categories': [], 'series': []}

            # Categories
            if hasattr(chart, 'plots') and len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, 'categories'):
                    chart_data['categories'] = [str(cat) for cat in plot.categories]

            # Series
            for series in chart.series:
                series_data = {
                    'name': series.name or 'Series',
                    'values': [v if v is not None else 0 for v in series.values]
                }
                chart_data['series'].append(series_data)

            return chart_data if chart_data['series'] else None

        except:
            return None

    def _render_bar_chart(self, ax, categories, series_list, chart_type):
        """Render bar/column chart"""
        x = np.arange(len(categories))
        width = 0.8 / len(series_list) if len(series_list) > 0 else 0.8

        for i, series in enumerate(series_list):
            offset = (i - len(series_list)/2) * width + width/2
            ax.bar(x + offset, series['values'], width, label=series['name'])

        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        ax.legend()
        ax.grid(axis='y', alpha=0.3)

    def _render_stacked_bar_chart(self, ax, categories, series_list, chart_type):
        """Render stacked bar chart"""
        x = np.arange(len(categories))
        width = 0.6

        bottom = np.zeros(len(categories))

        for series in series_list:
            values = np.array(series['values'])
            ax.bar(x, values, width, label=series['name'], bottom=bottom)
            bottom += values

        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        ax.legend()
        ax.grid(axis='y', alpha=0.3)

        if chart_type == XL_CHART_TYPE.COLUMN_STACKED_100:
            ax.set_ylabel('Percentage (%)')
            ax.set_ylim(0, 100)

    def _render_line_chart(self, ax, categories, series_list):
        """Render line chart"""
        x = np.arange(len(categories))

        for series in series_list:
            ax.plot(x, series['values'], marker='o', label=series['name'])

        ax.set_xticks(x)
        ax.set_xticklabels(categories, rotation=45, ha='right')
        ax.legend()
        ax.grid(alpha=0.3)

    def _render_pie_chart(self, ax, categories, series_list):
        """Render pie chart"""
        if len(series_list) > 0:
            series = series_list[0]
            ax.pie(series['values'], labels=categories, autopct='%1.1f%%', startangle=90)
            ax.axis('equal')
