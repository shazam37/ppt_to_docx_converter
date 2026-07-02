"""
Advanced test with images
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
from PIL import Image
from converter import PPTToDocConverter


def create_test_image(output_path, color, size=(200, 200)):
    """Create a simple colored test image"""
    img = Image.new('RGB', size, color=color)
    img.save(output_path)


def create_advanced_presentation(output_path):
    """Create presentation with images"""
    prs = Presentation()

    # Slide 1: Title with image
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_box.text_frame.text = "Advanced Test: Images & Mixed Content"

    # Create test images
    test_dir = Path(__file__).parent / "test_files"
    test_dir.mkdir(exist_ok=True)

    red_img = test_dir / "red.png"
    blue_img = test_dir / "blue.png"

    create_test_image(red_img, (255, 0, 0))
    create_test_image(blue_img, (0, 0, 255))

    # Add images
    left = Inches(1)
    top = Inches(2)
    slide1.shapes.add_picture(str(red_img), left, top, width=Inches(2))

    left = Inches(4)
    slide1.shapes.add_picture(str(blue_img), left, top, width=Inches(2))

    # Caption
    caption = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(6), Inches(0.5))
    caption.text_frame.text = "Red and Blue test images"

    # Slide 2: Mixed content
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Mixed Content Test"

    # Text box
    text_box = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(2))
    tf = text_box.text_frame
    tf.text = "This slide has both text and an image.\n\nThe image should be extracted and embedded in the Word document."

    # Image
    slide2.shapes.add_picture(str(red_img), Inches(5.5), Inches(2), width=Inches(1.5))

    prs.save(str(output_path))
    print(f"Advanced test presentation created: {output_path}")


def test_advanced_conversion():
    """Test advanced conversion"""
    test_dir = Path(__file__).parent / "test_files"
    test_dir.mkdir(exist_ok=True)

    input_file = test_dir / "test_advanced.pptx"
    output_file = test_dir / "test_advanced_output.docx"

    print("Creating advanced test presentation...")
    create_advanced_presentation(input_file)

    print("\nTesting advanced conversion...")

    def progress_callback(percent, message):
        print(f"Progress: {percent}% - {message}")

    def log_callback(message):
        print(f"Log: {message}")

    converter = PPTToDocConverter(
        progress_cb=progress_callback,
        log_cb=log_callback
    )

    try:
        converter.convert(input_file, output_file)
        print(f"\n✓ Advanced conversion successful!")
        print(f"Output file: {output_file}")
        print(f"File size: {output_file.stat().st_size / 1024:.2f} KB")

        # Verify images were extracted
        from docx import Document
        doc = Document(str(output_file))

        image_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_count += 1

        print(f"Images embedded: {image_count}")

        if image_count > 0:
            print("✓ Images successfully embedded in Word document")
            return True
        else:
            print("✗ No images found in Word document")
            return False

    except Exception as e:
        print(f"\n✗ Advanced conversion failed: {e}")
        import traceback
        traceback.print_exc()
        return False


if __name__ == "__main__":
    success = test_advanced_conversion()
    exit(0 if success else 1)
