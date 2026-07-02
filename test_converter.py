"""
Test script to create sample presentations and test the converter
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
from converter import PPTToDocConverter


def create_test_presentation(output_path):
    """Create a comprehensive test presentation with various elements"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = "Test Presentation"
    subtitle.text = "Comprehensive Feature Test for PPT to DOC Converter"

    # Slide 2: Text formatting
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Text Formatting Test"

    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    textbox = slide2.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    # Normal text
    p1 = text_frame.add_paragraph()
    p1.text = "This is normal text."
    p1.level = 0

    # Bold text
    p2 = text_frame.add_paragraph()
    p2.text = "This is bold text."
    p2.level = 0
    for run in p2.runs:
        run.font.bold = True

    # Italic text
    p3 = text_frame.add_paragraph()
    p3.text = "This is italic text."
    p3.level = 0
    for run in p3.runs:
        run.font.italic = True

    # Colored text
    p4 = text_frame.add_paragraph()
    p4.text = "This is colored text (red)."
    p4.level = 0
    for run in p4.runs:
        run.font.color.rgb = RGBColor(255, 0, 0)

    # Large text
    p5 = text_frame.add_paragraph()
    p5.text = "This is large text."
    p5.level = 0
    for run in p5.runs:
        run.font.size = Pt(24)

    # Slide 3: Bullet points
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Bullet Points Test"

    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    textbox3 = slide3.shapes.add_textbox(left, top, width, height)
    text_frame3 = textbox3.text_frame

    bullets = [
        "First bullet point",
        "Second bullet point",
        "Third bullet point with nested items:",
    ]

    for bullet in bullets:
        p = text_frame3.add_paragraph()
        p.text = bullet
        p.level = 0

    # Nested bullets
    nested = ["Nested item 1", "Nested item 2"]
    for item in nested:
        p = text_frame3.add_paragraph()
        p.text = item
        p.level = 1

    # Slide 4: Table
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title4 = slide4.shapes.title
    title4.text = "Table Test"

    rows = 4
    cols = 3
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(3)

    table = slide4.shapes.add_table(rows, cols, left, top, width, height).table

    # Header row
    table.cell(0, 0).text = "Product"
    table.cell(0, 1).text = "Quantity"
    table.cell(0, 2).text = "Price"

    # Data rows
    table.cell(1, 0).text = "Widget A"
    table.cell(1, 1).text = "100"
    table.cell(1, 2).text = "$10.00"

    table.cell(2, 0).text = "Widget B"
    table.cell(2, 1).text = "200"
    table.cell(2, 2).text = "$15.00"

    table.cell(3, 0).text = "Widget C"
    table.cell(3, 1).text = "150"
    table.cell(3, 2).text = "$12.50"

    # Slide 5: Multiple text boxes
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    title5_box = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title5_box.text_frame.text = "Multiple Text Boxes Test"

    # Box 1
    box1 = slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1.5))
    box1.text_frame.text = "Text Box 1: This is the first text box on the left side of the slide."

    # Box 2
    box2 = slide5.shapes.add_textbox(Inches(5), Inches(2), Inches(3), Inches(1.5))
    box2.text_frame.text = "Text Box 2: This is the second text box on the right side of the slide."

    # Box 3
    box3 = slide5.shapes.add_textbox(Inches(1), Inches(4), Inches(7), Inches(1.5))
    box3.text_frame.text = "Text Box 3: This is a wider text box at the bottom spanning most of the slide width."

    # Save presentation
    prs.save(str(output_path))
    print(f"Test presentation created: {output_path}")


def test_conversion():
    """Test the conversion process"""
    test_dir = Path(__file__).parent / "test_files"
    test_dir.mkdir(exist_ok=True)

    input_file = test_dir / "test_presentation.pptx"
    output_file = test_dir / "test_output.docx"

    # Create test presentation
    print("Creating test presentation...")
    create_test_presentation(input_file)

    # Test conversion
    print("\nTesting conversion...")

    def progress_callback(percent, message):
        print(f"Progress: {percent}% - {message}")

    def log_callback(message):
        print(f"Log: {message}")

    converter = PPTToDocConverter(
        progress_cb=progress_callback,
        log_cb=log_callback
    )

    try:
        converter.convert(input_file, output_file)
        print(f"\n✓ Conversion successful!")
        print(f"Output file: {output_file}")
        print(f"File size: {output_file.stat().st_size / 1024:.2f} KB")

        # Verify the output file exists and has content
        if output_file.exists() and output_file.stat().st_size > 0:
            print("\n✓ Output file verified - contains data")
            return True
        else:
            print("\n✗ Output file verification failed")
            return False

    except Exception as e:
        print(f"\n✗ Conversion failed: {e}")
        import traceback
        traceback.print_exc()
        return False


if __name__ == "__main__":
    success = test_conversion()
    exit(0 if success else 1)
